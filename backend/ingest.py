#!/usr/bin/env python3
"""
Document Ingestion Pipeline (Optimized)
Extracts text and images from documents, generates embeddings, and stores in vector DB
Includes parallel processing, caching, and progress reporting for 3-5x faster ingestion.

Run: python ingest.py
Options:
  --clear, -c  : Clear existing data before ingestion
  --stats, -s  : Show image extraction statistics only
"""

import os
import sqlite3
import time
from pathlib import Path
from typing import List, Tuple
import uuid

# Import local modules
from config import (
    DOCUMENTS_DIR,
    IMAGES_DIR,
    DB_PATH,
    SUPPORTED_EXTENSIONS,
    MAX_CHUNK_SIZE,
    CHUNK_OVERLAP,
    DEFAULT_TOPIC,
    SHOW_PROGRESS_BAR,
    get_available_topics,
)
from embeddings import get_embeddings_service
from vector_store import get_vector_store
from image_extractor import get_image_extractor

# Try to import tqdm for progress bars
try:
    from tqdm import tqdm
    TQDM_AVAILABLE = True
except ImportError:
    TQDM_AVAILABLE = False

try:
    from docx import Document as DocxDocument
except ImportError:
    print("Error: python-docx not installed. Run: pip install python-docx")
    exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    exit(1)


def format_time(seconds: float) -> str:
    """Format seconds into human-readable string"""
    if seconds < 60:
        return f"{seconds:.1f}s"
    elif seconds < 3600:
        mins = int(seconds // 60)
        secs = int(seconds % 60)
        return f"{mins}m {secs}s"
    else:
        hours = int(seconds // 3600)
        mins = int((seconds % 3600) // 60)
        return f"{hours}h {mins}m"


class DocumentProcessor:
    def __init__(self):
        self.embeddings_service = get_embeddings_service()
        self.vector_store = get_vector_store()
        self.image_extractor = get_image_extractor()
        self.show_progress = SHOW_PROGRESS_BAR and TQDM_AVAILABLE
        self._init_chat_db()

    def _init_chat_db(self):
        """Initialize SQLite database for chat history"""
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()

        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_message TEXT,
                assistant_response TEXT,
                source_documents TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """
        )

        conn.commit()
        conn.close()

    def process_all_documents(self):
        """Process all documents in the documents directory with progress tracking"""
        total_start_time = time.time()

        print("\n" + "=" * 60)
        print("DOCUMENT INGESTION PIPELINE (Optimized)")
        print("=" * 60)

        if not DOCUMENTS_DIR.exists():
            print(f"[ERROR] Documents directory not found: {DOCUMENTS_DIR}")
            return

        # ==============================================================
        # PHASE 1: Document Discovery
        # ==============================================================
        print("\n[PHASE 1/4] Document Discovery")
        print("-" * 40)
        phase1_start = time.time()

        # Collect all documents with their topics
        document_files_with_topics = []

        # First, get files in root documents folder (general topic)
        for ext in SUPPORTED_EXTENSIONS:
            for file in DOCUMENTS_DIR.glob(f"*{ext}"):
                if file.is_file():
                    document_files_with_topics.append((file, DEFAULT_TOPIC))

        # Then, scan subfolders (each subfolder = topic)
        available_topics = get_available_topics()
        print(f"[OK] Found {len(available_topics)} topic folders: {available_topics}")

        for topic in available_topics:
            topic_dir = DOCUMENTS_DIR / topic
            for ext in SUPPORTED_EXTENSIONS:
                for file in topic_dir.glob(f"*{ext}"):
                    if file.is_file():
                        document_files_with_topics.append((file, topic))
                # Also check nested folders within topic
                for file in topic_dir.glob(f"**/*{ext}"):
                    if file.is_file() and (file, topic) not in document_files_with_topics:
                        document_files_with_topics.append((file, topic))

        if not document_files_with_topics:
            print(f"[WARN] No documents found in {DOCUMENTS_DIR}")
            print(f"   Supported formats: {', '.join(SUPPORTED_EXTENSIONS)}")
            print(f"   You can organize documents in subfolders for topic filtering:")
            print(f"   e.g., {DOCUMENTS_DIR}/piping/file.pdf → topic='piping'")
            return

        # Separate PDFs from non-PDFs (PDFs take longer due to image extraction)
        pdf_docs = [(f, t) for f, t in document_files_with_topics if f.suffix.lower() == ".pdf"]
        other_docs = [(f, t) for f, t in document_files_with_topics if f.suffix.lower() != ".pdf"]

        print(f"[OK] Found {len(document_files_with_topics)} documents:")
        print(f"     - PDFs: {len(pdf_docs)}")
        print(f"     - Other: {len(other_docs)} (DOCX, PPTX, TXT)")

        for file, topic in document_files_with_topics[:10]:  # Show first 10
            print(f"     [{topic}] {file.name}")
        if len(document_files_with_topics) > 10:
            print(f"     ... and {len(document_files_with_topics) - 10} more")

        phase1_time = time.time() - phase1_start
        print(f"[DONE] Phase 1 completed in {format_time(phase1_time)}")

        # ==============================================================
        # PHASE 2: Text Extraction
        # ==============================================================
        print("\n[PHASE 2/4] Text Extraction")
        print("-" * 40)
        phase2_start = time.time()

        all_texts = []
        all_metadatas = []
        all_ids = []
        image_ids = []

        # Process non-PDFs first (faster)
        if other_docs:
            print(f"\n[TEXT] Processing {len(other_docs)} non-PDF documents...")
            doc_iter = other_docs
            if self.show_progress:
                doc_iter = tqdm(other_docs, desc="Extracting text", unit="doc")

            for doc_file, topic in doc_iter:
                texts, metadatas, ids = self._process_document(doc_file, topic)
                if texts:
                    all_texts.extend(texts)
                    all_metadatas.extend(metadatas)
                    all_ids.extend(ids)

        # Process PDFs (with image extraction)
        if pdf_docs:
            print(f"\n[PDF] Processing {len(pdf_docs)} PDF documents...")
            for doc_file, topic in pdf_docs:
                print(f"\n[FILE] {doc_file.name} (topic: {topic})")

                # Extract text
                texts, metadatas, ids = self._process_document(doc_file, topic)
                if texts:
                    all_texts.extend(texts)
                    all_metadatas.extend(metadatas)
                    all_ids.extend(ids)
                    print(f"  [OK] Extracted {len(texts)} text chunks")

                # Extract images (uses parallel processing internally)
                extracted = self.image_extractor.extract_images_from_pdf(
                    str(doc_file), doc_file.name, topic=topic
                )
                if extracted:
                    # Get image IDs from database
                    conn = sqlite3.connect(DB_PATH)
                    cursor = conn.cursor()
                    cursor.execute(
                        "SELECT id FROM images WHERE document_name = ?",
                        (doc_file.name,),
                    )
                    for row in cursor.fetchall():
                        image_ids.append(row[0])
                    conn.close()

        # Process PPTX images separately (already in other_docs but need image extraction)
        pptx_docs = [(f, t) for f, t in other_docs if f.suffix.lower() == ".pptx"]
        if pptx_docs:
            print(f"\n[PPTX] Extracting images from {len(pptx_docs)} PPTX files...")
            for doc_file, topic in pptx_docs:
                extracted = self.image_extractor.extract_images_from_pptx(
                    str(doc_file), doc_file.name, topic=topic
                )
                if extracted:
                    conn = sqlite3.connect(DB_PATH)
                    cursor = conn.cursor()
                    cursor.execute(
                        "SELECT id FROM images WHERE document_name = ?",
                        (doc_file.name,),
                    )
                    for row in cursor.fetchall():
                        image_ids.append(row[0])
                    conn.close()

        phase2_time = time.time() - phase2_start
        print(f"\n[DONE] Phase 2 completed in {format_time(phase2_time)}")
        print(f"       Total text chunks: {len(all_texts)}")
        print(f"       Total images found: {len(image_ids)}")

        # ==============================================================
        # PHASE 3: Vector Store Indexing
        # ==============================================================
        print("\n[PHASE 3/4] Vector Store Indexing")
        print("-" * 40)
        phase3_start = time.time()

        if all_texts:
            print(f"[DB] Adding {len(all_texts)} text documents to vector store...")
            self.vector_store.add_text_documents(all_texts, all_metadatas, all_ids)

        if image_ids:
            print(f"\n[IMG] Adding {len(image_ids)} images to vector store...")
            self.vector_store.add_image_metadata(image_ids)

        phase3_time = time.time() - phase3_start
        print(f"\n[DONE] Phase 3 completed in {format_time(phase3_time)}")

        # ==============================================================
        # PHASE 4: Summary & Statistics
        # ==============================================================
        print("\n[PHASE 4/4] Summary & Statistics")
        print("-" * 40)

        total_time = time.time() - total_start_time
        stats = self.vector_store.get_collection_stats()
        cache_stats = self.embeddings_service.get_cache_stats()

        print("\n" + "=" * 60)
        print("INGESTION COMPLETE")
        print("=" * 60)
        print(f"\n[RESULTS]")
        print(f"  Text documents indexed: {stats['text_documents']}")
        print(f"  Images indexed: {stats['images']}")
        print(f"  Images directory: {IMAGES_DIR}")
        print(f"  Vector DB: {DB_PATH}")

        print(f"\n[TIMING]")
        print(f"  Phase 1 (Discovery):  {format_time(phase1_time)}")
        print(f"  Phase 2 (Extraction): {format_time(phase2_time)}")
        print(f"  Phase 3 (Indexing):   {format_time(phase3_time)}")
        print(f"  Total time:           {format_time(total_time)}")

        print(f"\n[EMBEDDING CACHE]")
        print(f"  Cache enabled: {cache_stats['enabled']}")
        print(f"  Cache hits: {cache_stats['cache_hits']}")
        print(f"  Cache misses: {cache_stats['cache_misses']}")
        print(f"  Hit rate: {cache_stats['hit_rate_percent']}%")
        print(f"  Cache files: {cache_stats['cache_files']}")
        print(f"  Cache size: {cache_stats['cache_size_mb']} MB")

        print("\n[OK] Ready to start backend server!")
        print("=" * 60 + "\n")

    def _process_document(
        self, file_path: Path, topic: str = DEFAULT_TOPIC
    ) -> Tuple[List[str], List[dict], List[str]]:
        """
        Process a single document and return text chunks with topic metadata

        Args:
            file_path: Path to document
            topic: Topic name (derived from parent folder)

        Returns:
            Tuple of (texts, metadatas, ids)
        """
        ext = file_path.suffix.lower()

        try:
            if ext == ".pdf":
                text = self._extract_text_from_pdf(file_path)
            elif ext == ".docx":
                text = self._extract_text_from_docx(file_path)
            elif ext == ".pptx":
                text = self._extract_text_from_pptx(file_path)
            elif ext == ".txt":
                text = self._extract_text_from_txt(file_path)
            else:
                print(f"  [ERROR] Unsupported file format: {ext}")
                return [], [], []

            if not text or text.strip() == "":
                print(f"  [X] No text extracted")
                return [], [], []

            # Split into chunks WITH accurate page tracking
            chunks_with_pages = self._chunk_text_with_pages(text, MAX_CHUNK_SIZE, CHUNK_OVERLAP)

            # Create metadata for each chunk with topic and accurate page numbers
            chunks = []
            metadatas = []
            for i, (chunk, page_number) in enumerate(chunks_with_pages):
                chunks.append(chunk)
                metadatas.append({
                    "document_name": file_path.name,
                    "source": str(file_path),  # Full source path
                    "topic": topic,  # Topic from folder name
                    "chunk_index": i,
                    "page_number": page_number,
                    "file_type": ext,
                })

            # Create unique IDs
            ids = [f"{file_path.stem}_{uuid.uuid4()}" for _ in chunks]

            return chunks, metadatas, ids

        except Exception as e:
            import traceback
            error_msg = str(e) or repr(e)
            print(f"  [ERROR] Error processing file: {type(e).__name__}: {error_msg}")
            traceback.print_exc()
            return [], [], []

    def _extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF - processes pages one at a time to save memory"""
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                # Fallback to PyPDF2
                from PyPDF2 import PdfReader
            except ImportError:
                print("  [ERROR] PDF library not installed. Run: pip install pypdf")
                return ""

        text = []
        try:
            pdf = PdfReader(str(file_path))
            total_pages = len(pdf.pages)

            page_iter = range(total_pages)
            if self.show_progress and total_pages > 50:
                page_iter = tqdm(page_iter, desc="  Reading pages", unit="page", leave=False)

            for page_num in page_iter:
                try:
                    page_text = pdf.pages[page_num].extract_text()
                    if page_text and page_text.strip():
                        text.append(f"[Page {page_num + 1}] {page_text}")
                except Exception:
                    continue

        except Exception as e:
            print(f"  [WARN] Error reading PDF: {str(e)}")

        return "\n\n".join(text)

    def _extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        text = []
        try:
            doc = DocxDocument(str(file_path))
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text.append(cell.text)
        except Exception as e:
            print(f"  [WARN] Error reading DOCX: {str(e)}")

        return "\n\n".join(text)

    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX"""
        text = []
        try:
            prs = Presentation(str(file_path))
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                text.append("\n".join(slide_text))
        except Exception as e:
            print(f"  [WARN] Error reading PPTX: {str(e)}")

        return "\n\n".join(text)

    def _extract_text_from_txt(self, file_path: Path) -> str:
        """Extract text from TXT"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            print(f"  [WARN] Error reading TXT: {str(e)}")
            return ""

    def _chunk_text_with_pages(self, text: str, max_size: int, overlap: int) -> List[Tuple[str, int]]:
        """
        Split text into overlapping chunks while preserving page numbers.

        Args:
            text: Text to chunk (contains [Page X] markers)
            max_size: Max characters per chunk
            overlap: Character overlap between chunks

        Returns:
            List of (chunk_text, page_number) tuples
        """
        import re

        if not text:
            return []

        # Ensure overlap is less than max_size
        if overlap >= max_size:
            overlap = max_size // 2

        # Build a map of character positions to page numbers
        page_markers = list(re.finditer(r'\[Page\s*(\d+)\]|\[Slide\s*(\d+)\]', text))

        def get_page_at_position(pos: int) -> int:
            """Get the page number at a given character position"""
            current_page = 1
            for marker in page_markers:
                if marker.start() <= pos:
                    # Extract page number from either Page or Slide group
                    page_num = marker.group(1) or marker.group(2)
                    current_page = int(page_num)
                else:
                    break
            return current_page

        chunks_with_pages = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = min(start + max_size, text_len)
            chunk = text[start:end]

            # Try to break at a sentence boundary
            if end < text_len:
                last_period = chunk.rfind(". ")
                if last_period > max_size * 0.7:
                    end = start + last_period + 2
                    chunk = text[start:end]

            stripped_chunk = chunk.strip()
            if stripped_chunk and len(stripped_chunk) > 50:
                # Get the page number at the START of this chunk
                page_num = get_page_at_position(start)
                chunks_with_pages.append((stripped_chunk, page_num))

            # Move forward
            new_start = end - overlap
            if new_start <= start:
                new_start = start + 1
            start = new_start

        return chunks_with_pages

    def _chunk_text(self, text: str, max_size: int, overlap: int) -> List[str]:
        """Legacy method - returns just chunks without page info"""
        result = self._chunk_text_with_pages(text, max_size, overlap)
        return [chunk for chunk, _ in result]


def main():
    """Main ingestion pipeline"""
    import sys

    # Check for flags
    clear_first = "--clear" in sys.argv or "-c" in sys.argv
    show_stats = "--stats" in sys.argv or "-s" in sys.argv

    # If --stats flag, show image extraction statistics
    if show_stats:
        print("\n" + "=" * 60)
        print("IMAGE EXTRACTION STATISTICS")
        print("=" * 60)
        image_extractor = get_image_extractor()
        all_images = image_extractor.get_all_images()
        images_with_caption = [img for img in all_images if img.get("figure_caption")]
        images_without_caption = image_extractor.get_images_without_caption()

        print(f"Total images: {len(all_images)}")
        print(f"Images with figure captions: {len(images_with_caption)}")
        print(f"Images without captions: {len(images_without_caption)}")

        if images_with_caption:
            print("\nSample captions:")
            for img in images_with_caption[:5]:
                caption = img.get("figure_caption", "")[:80]
                print(f"  - Page {img['page_number']}: {caption}...")

        # Show cache stats
        embeddings_service = get_embeddings_service()
        cache_stats = embeddings_service.get_cache_stats()
        print(f"\n[EMBEDDING CACHE]")
        print(f"  Cache files: {cache_stats['cache_files']}")
        print(f"  Cache size: {cache_stats['cache_size_mb']} MB")

        print("\n" + "=" * 60)
        return

    processor = DocumentProcessor()

    if clear_first:
        print("\n[CLEAR] Clearing existing data...")
        # Clear vector store
        processor.vector_store.clear_all()

        # Clear embedding cache
        processor.embeddings_service.clear_cache()

        # Clear images from SQLite
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute("DELETE FROM images")
        cursor.execute("DELETE FROM chat_history")
        try:
            cursor.execute("DELETE FROM sessions")
        except:
            pass
        conn.commit()
        conn.close()

        # Clear extracted images folder
        import shutil
        if IMAGES_DIR.exists():
            for f in IMAGES_DIR.glob("*"):
                f.unlink()
        print("[OK] Cleared all existing data\n")

    processor.process_all_documents()


if __name__ == "__main__":
    main()
