import os
import sqlite3
import base64
import requests
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import hashlib
from config import IMAGES_DIR, DB_PATH, OLLAMA_API_URL, OLLAMA_TIMEOUT, DEFAULT_TOPIC

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("[WARN] PyMuPDF not installed. Run: pip install PyMuPDF")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("[WARN] python-pptx not installed. Run: pip install python-pptx")


# =============================================================================
# CONFIGURATION FOR LAYOUT-AWARE EXTRACTION
# =============================================================================

# Figure caption patterns (case-insensitive)
FIGURE_CAPTION_PATTERNS = [
    r'(Fig\.?\s*\d+[-.]?\d*[.:]\s*[^\n]{5,150})',
    r'(Figure\s+\d+[-.]?\d*[.:]\s*[^\n]{5,150})',
    r'(FIGURE\s+\d+[-.]?\d*[.:]\s*[^\n]{5,150})',
    r'(Fig\.?\s*\d+[-.]?\d*\s*[-–—]\s*[^\n]{5,150})',
    r'(Figure\s+\d+[-.]?\d*\s*[-–—]\s*[^\n]{5,150})',
]

# Vertical distance threshold for nearby text (in points, 1 point = 1/72 inch)
VERTICAL_PROXIMITY_THRESHOLD = 50  # ~0.7 inches

# Section heading patterns
SECTION_HEADING_PATTERNS = [
    r'^[A-Z][A-Za-z\s]+$',  # Title case headers
    r'^\d+\.?\s+[A-Z][A-Za-z\s]+',  # Numbered headers
    r'^[A-Z\s]+$',  # ALL CAPS headers
]


class ImageExtractor:
    def __init__(self):
        self.images_dir = IMAGES_DIR
        self._init_image_db()

    def _init_image_db(self):
        """Initialize SQLite database for image metadata"""
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS images (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                image_hash TEXT UNIQUE,
                image_path TEXT,
                document_name TEXT,
                page_number INTEGER,
                context_text TEXT,
                image_description TEXT,
                figure_caption TEXT,
                bbox_x0 REAL,
                bbox_y0 REAL,
                bbox_x1 REAL,
                bbox_y1 REAL,
                embedding BLOB,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        # Add columns if they don't exist (for migration)
        migration_columns = [
            ('context_text', 'TEXT'),
            ('image_description', 'TEXT'),
            ('figure_caption', 'TEXT'),
            ('bbox_x0', 'REAL'),
            ('bbox_y0', 'REAL'),
            ('bbox_x1', 'REAL'),
            ('bbox_y1', 'REAL'),
            ('topic', 'TEXT'),  # Topic from folder name
        ]
        for col_name, col_type in migration_columns:
            try:
                cursor.execute(f"ALTER TABLE images ADD COLUMN {col_name} {col_type}")
            except:
                pass  # Column already exists
        conn.commit()
        conn.close()

    # =========================================================================
    # LAYOUT-AWARE TEXT EXTRACTION
    # =========================================================================

    def _extract_text_blocks_with_coords(self, page) -> List[Dict]:
        """
        Extract all text blocks from a page with their bounding box coordinates.

        Args:
            page: PyMuPDF page object

        Returns:
            List of dicts with text and bbox info
        """
        text_blocks = []

        # Get text blocks with their positions
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

        for block in blocks:
            if block.get("type") == 0:  # Text block
                bbox = block.get("bbox", (0, 0, 0, 0))

                # Combine all text in the block
                block_text = ""
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        block_text += span.get("text", "")
                    block_text += " "

                block_text = block_text.strip()
                if block_text:
                    text_blocks.append({
                        "text": block_text,
                        "bbox": bbox,  # (x0, y0, x1, y1)
                        "x0": bbox[0],
                        "y0": bbox[1],
                        "x1": bbox[2],
                        "y1": bbox[3],
                        "center_y": (bbox[1] + bbox[3]) / 2,
                        "center_x": (bbox[0] + bbox[2]) / 2,
                    })

        return text_blocks

    def _detect_figure_caption(self, text_blocks: List[Dict]) -> List[Dict]:
        """
        Detect figure captions from text blocks using common patterns.

        Args:
            text_blocks: List of text blocks with coordinates

        Returns:
            List of caption dicts with text and bbox
        """
        captions = []

        for block in text_blocks:
            text = block["text"]

            # Check against all caption patterns
            for pattern in FIGURE_CAPTION_PATTERNS:
                matches = re.findall(pattern, text, re.IGNORECASE)
                if matches:
                    for match in matches:
                        captions.append({
                            "text": match.strip(),
                            "bbox": block["bbox"],
                            "y0": block["y0"],
                            "y1": block["y1"],
                            "center_y": block["center_y"],
                            "full_block": block,
                        })

        return captions

    def _find_nearby_text_for_image(self, image_bbox: Tuple[float, float, float, float],
                                     text_blocks: List[Dict],
                                     captions: List[Dict],
                                     vertical_threshold: float = VERTICAL_PROXIMITY_THRESHOLD) -> Dict:
        """
        Find the most relevant text context for an image based on spatial layout.

        Priority order:
        1. Explicit figure captions (above, below, or beside)
        2. Section headings
        3. Closest descriptive paragraphs

        Args:
            image_bbox: (x0, y0, x1, y1) of the image
            text_blocks: All text blocks on the page
            captions: Detected figure captions
            vertical_threshold: Max distance to consider text as "nearby"

        Returns:
            Dict with 'figure_caption', 'context_text', 'section_heading'
        """
        img_x0, img_y0, img_x1, img_y1 = image_bbox
        img_center_y = (img_y0 + img_y1) / 2
        img_center_x = (img_x0 + img_x1) / 2

        result = {
            "figure_caption": None,
            "section_heading": None,
            "context_text": "",
            "nearby_text_blocks": [],
        }

        # PRIORITY 1: Find explicit figure caption
        best_caption = None
        best_caption_distance = float('inf')

        for caption in captions:
            # Check if caption is near the image (above or below)
            cap_center_y = caption["center_y"]

            # Caption typically appears above or below the image
            distance = min(
                abs(cap_center_y - img_y0),  # Distance to top of image
                abs(cap_center_y - img_y1),  # Distance to bottom of image
            )

            # Also check horizontal overlap
            cap_bbox = caption["bbox"]
            horizontal_overlap = (
                cap_bbox[0] < img_x1 and cap_bbox[2] > img_x0
            )

            if distance < vertical_threshold * 2 and horizontal_overlap:
                if distance < best_caption_distance:
                    best_caption_distance = distance
                    best_caption = caption["text"]

        if best_caption:
            result["figure_caption"] = best_caption

        # PRIORITY 2: Find section headings above the image
        for block in text_blocks:
            if block["y1"] < img_y0:  # Above the image
                text = block["text"].strip()
                for pattern in SECTION_HEADING_PATTERNS:
                    if re.match(pattern, text) and len(text) < 100:
                        # Check if this is the closest heading
                        distance = img_y0 - block["y1"]
                        if distance < vertical_threshold * 3:
                            result["section_heading"] = text
                            break

        # PRIORITY 3: Collect nearby text blocks (above and below)
        nearby_blocks = []

        for block in text_blocks:
            block_center_y = block["center_y"]

            # Check if block is vertically close to the image
            is_above = block["y1"] <= img_y0 and (img_y0 - block["y1"]) < vertical_threshold
            is_below = block["y0"] >= img_y1 and (block["y0"] - img_y1) < vertical_threshold

            # Check horizontal overlap (at least 30%)
            block_x0, block_x1 = block["x0"], block["x1"]
            overlap_start = max(block_x0, img_x0)
            overlap_end = min(block_x1, img_x1)
            horizontal_overlap = (overlap_end - overlap_start) / max(1, img_x1 - img_x0)

            if (is_above or is_below) and horizontal_overlap > 0.3:
                # Calculate distance for sorting
                if is_above:
                    distance = img_y0 - block["y1"]
                else:
                    distance = block["y0"] - img_y1

                nearby_blocks.append({
                    "text": block["text"],
                    "distance": distance,
                    "position": "above" if is_above else "below",
                })

        # Sort by distance and combine
        nearby_blocks.sort(key=lambda x: x["distance"])
        result["nearby_text_blocks"] = nearby_blocks[:5]  # Top 5 closest blocks

        # Build context string (combine caption + heading + nearby text)
        context_parts = []

        if result["figure_caption"]:
            context_parts.append(f"CAPTION: {result['figure_caption']}")

        if result["section_heading"]:
            context_parts.append(f"SECTION: {result['section_heading']}")

        for block in nearby_blocks[:3]:
            if len(block["text"]) > 20:  # Skip very short text
                context_parts.append(block["text"])

        result["context_text"] = " | ".join(context_parts) if context_parts else ""

        return result

    def _get_image_bbox_on_page(self, page, xref: int) -> Optional[Tuple[float, float, float, float]]:
        """
        Get the bounding box of an image on a page.

        Args:
            page: PyMuPDF page object
            xref: Image xref number

        Returns:
            Bounding box (x0, y0, x1, y1) or None
        """
        try:
            # Get all image instances on the page
            image_list = page.get_images(full=True)

            # Find the image with matching xref
            for img_info in image_list:
                if img_info[0] == xref:
                    # Get the image rectangle
                    for img_rect in page.get_image_rects(xref):
                        return (img_rect.x0, img_rect.y0, img_rect.x1, img_rect.y1)

            return None
        except Exception:
            return None

    # =========================================================================
    # PDF IMAGE EXTRACTION (LAYOUT-AWARE)
    # =========================================================================

    def extract_images_from_pdf(self, pdf_path: str, doc_name: str, topic: str = None) -> List[Dict]:
        """
        Extract figures/diagrams from PDF with layout-aware context extraction.

        For each image:
        - Extract bounding box coordinates
        - Extract all text blocks with spatial coordinates
        - Detect figure captions using patterns
        - Find spatially-related text (above, below, beside)

        Args:
            pdf_path: Path to PDF file
            doc_name: Name of document for metadata
            topic: Topic name (from folder)

        Returns:
            List of extracted image info dicts
        """
        topic = topic or DEFAULT_TOPIC
        if not PYMUPDF_AVAILABLE:
            print("  [WARN] PyMuPDF not available, skipping image extraction")
            return []

        extracted_images = []

        try:
            pdf = fitz.open(pdf_path)
            total_pages = len(pdf)
            print(f"  [IMG] Extracting figures from {total_pages} pages (layout-aware)...")

            image_count = 0

            for page_num in range(total_pages):
                try:
                    page = pdf[page_num]

                    # Step 1: Extract all text blocks with coordinates
                    text_blocks = self._extract_text_blocks_with_coords(page)

                    # Step 2: Detect figure captions from text blocks
                    captions = self._detect_figure_caption(text_blocks)

                    # Step 3: Get all images on this page
                    image_list = page.get_images(full=True)

                    for img_index, img_info in enumerate(image_list):
                        try:
                            xref = img_info[0]

                            # Extract the image
                            base_image = pdf.extract_image(xref)
                            if not base_image:
                                continue

                            image_bytes = base_image["image"]
                            image_ext = base_image.get("ext", "png")
                            width = base_image.get("width", 0)
                            height = base_image.get("height", 0)

                            # Filter: Skip tiny images (icons, bullets, decorations)
                            if width < 150 or height < 150:
                                continue

                            # Filter: Skip very narrow/tall images (likely borders)
                            aspect_ratio = width / height if height > 0 else 0
                            if aspect_ratio > 5 or aspect_ratio < 0.2:
                                continue

                            # Filter: Skip very large images (likely backgrounds)
                            if width > 2500 or height > 2500:
                                continue

                            # Filter: Skip very small file sizes (simple graphics)
                            if len(image_bytes) < 5000:
                                continue

                            # Generate hash for deduplication
                            image_hash = hashlib.md5(image_bytes).hexdigest()

                            # Check if already exists
                            conn = sqlite3.connect(DB_PATH)
                            cursor = conn.cursor()
                            cursor.execute("SELECT id FROM images WHERE image_hash = ?", (image_hash,))
                            if cursor.fetchone():
                                conn.close()
                                continue
                            conn.close()

                            # Step 4: Get image bounding box
                            image_bbox = self._get_image_bbox_on_page(page, xref)

                            # Step 5: Find layout-aware context
                            if image_bbox:
                                context_info = self._find_nearby_text_for_image(
                                    image_bbox, text_blocks, captions
                                )
                            else:
                                # Fallback: use page-level caption detection
                                context_info = {
                                    "figure_caption": captions[0]["text"] if captions else None,
                                    "context_text": " | ".join([c["text"] for c in captions[:2]]) if captions else "",
                                    "section_heading": None,
                                }
                                image_bbox = (0, 0, 0, 0)  # Unknown bbox

                            # Save image
                            image_filename = f"{Path(pdf_path).stem}_p{page_num + 1}_fig{img_index}.{image_ext}"
                            image_path = self.images_dir / image_filename

                            with open(image_path, "wb") as f:
                                f.write(image_bytes)

                            # Store metadata with layout-aware context and topic
                            image_id = self._store_image_metadata(
                                image_hash=image_hash,
                                image_path=str(image_path),
                                document_name=doc_name,
                                page_number=page_num + 1,
                                context_text=context_info.get("context_text", ""),
                                figure_caption=context_info.get("figure_caption"),
                                bbox=image_bbox,
                                topic=topic
                            )

                            if image_id:
                                extracted_images.append({
                                    "id": image_id,
                                    "path": str(image_path),
                                    "page": page_num + 1,
                                    "caption": context_info.get("figure_caption"),
                                })
                                image_count += 1

                        except Exception as img_err:
                            continue

                except Exception as page_err:
                    continue

                if (page_num + 1) % 100 == 0:
                    print(f"  [OK] Processed {page_num + 1}/{total_pages} pages, found {image_count} figures...")

            pdf.close()
            print(f"  [OK] Extracted {len(extracted_images)} figures/diagrams with layout context")

        except Exception as e:
            print(f"  [ERROR] Error extracting images: {str(e)}")

        return extracted_images

    # =========================================================================
    # PPTX IMAGE EXTRACTION (LAYOUT-AWARE)
    # =========================================================================

    def extract_images_from_pptx(self, pptx_path: str, doc_name: str, topic: str = None) -> List[Dict]:
        """
        Extract images from PowerPoint with layout-aware context.

        Args:
            pptx_path: Path to PPTX file
            doc_name: Name of document for metadata
            topic: Topic name (from folder)

        Returns:
            List of extracted image info dicts
        """
        topic = topic or DEFAULT_TOPIC
        if not PPTX_AVAILABLE:
            print("  [WARN] python-pptx not available, skipping PPTX image extraction")
            return []

        extracted_images = []

        try:
            prs = Presentation(pptx_path)
            total_slides = len(prs.slides)
            print(f"  [IMG] Extracting images from {total_slides} slides (layout-aware)...")

            image_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    # Collect all text shapes with positions
                    text_shapes = []
                    image_shapes = []

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_shapes.append({
                                "text": shape.text.strip(),
                                "top": shape.top if hasattr(shape, "top") else 0,
                                "left": shape.left if hasattr(shape, "left") else 0,
                                "width": shape.width if hasattr(shape, "width") else 0,
                                "height": shape.height if hasattr(shape, "height") else 0,
                            })

                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_shapes.append(shape)

                    # Detect figure captions from text shapes
                    captions = []
                    for ts in text_shapes:
                        for pattern in FIGURE_CAPTION_PATTERNS:
                            if re.search(pattern, ts["text"], re.IGNORECASE):
                                captions.append(ts)
                                break

                    # Build slide context (title + main text)
                    slide_context_parts = []
                    for ts in text_shapes[:5]:
                        if len(ts["text"]) > 10:
                            slide_context_parts.append(ts["text"])
                    slide_context = " | ".join(slide_context_parts) if slide_context_parts else f"Slide {slide_num}"

                    # Process each image shape
                    for shape_idx, shape in enumerate(image_shapes):
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext

                            # Filter small images
                            if shape.width.emu < 500000 or shape.height.emu < 500000:
                                continue

                            if len(image_bytes) < 3000:
                                continue

                            # Generate hash for deduplication
                            image_hash = hashlib.md5(image_bytes).hexdigest()

                            # Check if already exists
                            conn = sqlite3.connect(DB_PATH)
                            cursor = conn.cursor()
                            cursor.execute("SELECT id FROM images WHERE image_hash = ?", (image_hash,))
                            if cursor.fetchone():
                                conn.close()
                                continue
                            conn.close()

                            # Find nearest caption
                            figure_caption = None
                            if captions:
                                shape_center_y = shape.top + shape.height / 2 if hasattr(shape, "top") else 0
                                best_caption = None
                                best_distance = float('inf')

                                for cap in captions:
                                    cap_center_y = cap["top"] + cap["height"] / 2
                                    distance = abs(cap_center_y - shape_center_y)
                                    if distance < best_distance:
                                        best_distance = distance
                                        best_caption = cap["text"]

                                if best_caption:
                                    figure_caption = best_caption

                            # Build context
                            context_parts = []
                            if figure_caption:
                                context_parts.append(f"CAPTION: {figure_caption}")
                            context_parts.append(slide_context)
                            context_text = " | ".join(context_parts)

                            # Save image
                            image_filename = f"{Path(pptx_path).stem}_slide{slide_num}_img{shape_idx}.{image_ext}"
                            image_path = self.images_dir / image_filename

                            with open(image_path, "wb") as f:
                                f.write(image_bytes)

                            # Get bbox (convert EMU to approximate points)
                            bbox = (
                                shape.left / 914400 * 72 if hasattr(shape, "left") else 0,
                                shape.top / 914400 * 72 if hasattr(shape, "top") else 0,
                                (shape.left + shape.width) / 914400 * 72 if hasattr(shape, "left") else 0,
                                (shape.top + shape.height) / 914400 * 72 if hasattr(shape, "top") else 0,
                            )

                            # Store metadata with topic
                            image_id = self._store_image_metadata(
                                image_hash=image_hash,
                                image_path=str(image_path),
                                document_name=doc_name,
                                page_number=slide_num,
                                context_text=context_text,
                                figure_caption=figure_caption,
                                bbox=bbox,
                                topic=topic
                            )

                            if image_id:
                                extracted_images.append({
                                    "id": image_id,
                                    "path": str(image_path),
                                    "page": slide_num,
                                    "caption": figure_caption,
                                })
                                image_count += 1

                        except Exception:
                            continue

                    # Also check group shapes
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            for group_shape in shape.shapes:
                                if hasattr(group_shape, 'shape_type') and group_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    try:
                                        image = group_shape.image
                                        image_bytes = image.blob
                                        image_ext = image.ext

                                        if len(image_bytes) < 3000:
                                            continue

                                        image_hash = hashlib.md5(image_bytes).hexdigest()

                                        conn = sqlite3.connect(DB_PATH)
                                        cursor = conn.cursor()
                                        cursor.execute("SELECT id FROM images WHERE image_hash = ?", (image_hash,))
                                        if cursor.fetchone():
                                            conn.close()
                                            continue
                                        conn.close()

                                        image_filename = f"{Path(pptx_path).stem}_slide{slide_num}_grp{shape_idx}.{image_ext}"
                                        image_path = self.images_dir / image_filename

                                        with open(image_path, "wb") as f:
                                            f.write(image_bytes)

                                        image_id = self._store_image_metadata(
                                            image_hash=image_hash,
                                            image_path=str(image_path),
                                            document_name=doc_name,
                                            page_number=slide_num,
                                            context_text=slide_context,
                                            figure_caption=None,
                                            bbox=(0, 0, 0, 0),
                                            topic=topic
                                        )

                                        if image_id:
                                            extracted_images.append({
                                                "id": image_id,
                                                "path": str(image_path),
                                                "page": slide_num,
                                                "caption": None,
                                            })
                                            image_count += 1
                                    except Exception:
                                        continue

                except Exception:
                    continue

                if slide_num % 50 == 0:
                    print(f"  [OK] Processed {slide_num}/{total_slides} slides, found {image_count} images...")

            print(f"  [OK] Extracted {len(extracted_images)} images from PPTX with layout context")

        except Exception as e:
            print(f"  [ERROR] Error extracting PPTX images: {str(e)}")

        return extracted_images

    # =========================================================================
    # DATABASE OPERATIONS
    # =========================================================================

    def _store_image_metadata(self, image_hash: str, image_path: str,
                              document_name: str, page_number: int,
                              context_text: str = "",
                              figure_caption: str = None,
                              bbox: Tuple[float, float, float, float] = None,
                              topic: str = None) -> Optional[int]:
        """Store image metadata in SQLite with layout info and topic"""
        try:
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()

            bbox = bbox or (0, 0, 0, 0)
            topic = topic or DEFAULT_TOPIC

            cursor.execute("""
                INSERT OR IGNORE INTO images
                (image_hash, image_path, document_name, page_number, context_text,
                 figure_caption, bbox_x0, bbox_y0, bbox_x1, bbox_y1, topic)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (image_hash, image_path, document_name, page_number, context_text,
                  figure_caption, bbox[0], bbox[1], bbox[2], bbox[3], topic))
            conn.commit()
            image_id = cursor.lastrowid
            conn.close()
            return image_id
        except Exception as e:
            print(f"  [ERROR] Error storing image metadata: {str(e)}")
            return None

    def get_image_by_id(self, image_id: int) -> Optional[Dict]:
        """Retrieve image metadata by ID"""
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, image_path, document_name, page_number, figure_caption, context_text
            FROM images WHERE id = ?
        """, (image_id,))
        row = cursor.fetchone()
        conn.close()

        if row:
            return {
                "id": row[0],
                "image_path": row[1],
                "document_name": row[2],
                "page_number": row[3],
                "figure_caption": row[4],
                "context_text": row[5],
            }
        return None

    # =========================================================================
    # IMAGE RETRIEVAL (LAYOUT-AWARE, NO LLAVA)
    # =========================================================================

    def find_relevant_images(self, document_name: str, page_numbers: List[int],
                             search_terms: List[str] = None,
                             text_content: str = None,
                             max_images: int = 3,
                             query: str = None) -> List[Dict]:
        """
        Find relevant images using layout-derived context (NO LLaVA).

        Retrieval priority:
        1. Figure captions matching the query
        2. Layout-derived context text
        3. Page-based matching (as approximate hint only)

        If no related image exists for the topic, returns empty list.

        Args:
            document_name: Document to search in
            page_numbers: Hint pages (not truth, just hints)
            search_terms: Keywords from query
            text_content: Text content for figure reference detection
            max_images: Maximum images to return
            query: Original user query

        Returns:
            List of relevant images (empty if no good match)
        """
        found_images = []
        seen_ids = set()

        # Minimum relevance threshold - if no image scores above this, return nothing
        MIN_RELEVANCE_SCORE = 5

        # STRATEGY 1: Search by figure caption (HIGHEST PRIORITY)
        if query or search_terms:
            caption_matches = self._search_by_figure_caption(
                query or " ".join(search_terms or []),
                document_name,
                max_images * 2
            )

            for img in caption_matches:
                if img["id"] not in seen_ids and len(found_images) < max_images:
                    if img.get("score", 0) >= MIN_RELEVANCE_SCORE:
                        seen_ids.add(img["id"])
                        found_images.append(img)
                        print(f"  [IMG] Found by caption match: '{img.get('figure_caption', '')[:50]}...'")

        # STRATEGY 2: Search by layout-derived context
        if len(found_images) < max_images and (query or search_terms):
            context_matches = self._search_by_context(
                query or " ".join(search_terms or []),
                document_name,
                max_images * 2
            )

            for img in context_matches:
                if img["id"] not in seen_ids and len(found_images) < max_images:
                    if img.get("score", 0) >= MIN_RELEVANCE_SCORE:
                        seen_ids.add(img["id"])
                        found_images.append(img)
                        print(f"  [IMG] Found by context match: page {img['page_number']}")

        # STRATEGY 3: Figure reference in text (e.g., "see Figure 3-1")
        if text_content and len(found_images) < max_images:
            figure_refs = self._extract_figure_references(text_content)

            for fig_ref in figure_refs:
                if len(found_images) >= max_images:
                    break

                ref_matches = self._search_by_figure_reference(fig_ref, document_name)
                for img in ref_matches:
                    if img["id"] not in seen_ids and len(found_images) < max_images:
                        seen_ids.add(img["id"])
                        found_images.append(img)
                        print(f"  [IMG] Found by figure reference: {fig_ref}")

        # NOTE: Page-based matching removed as primary strategy
        # Page numbers are just hints, not truth. Only use if we have strong caption/context match.

        if not found_images:
            print("  [IMG] No relevant images found for query (no good caption/context match)")

        return found_images[:max_images]

    def _search_by_figure_caption(self, query: str, document_name: str = None, limit: int = 5) -> List[Dict]:
        """
        Search images by matching query against figure captions.

        Args:
            query: Search query
            document_name: Optional document filter
            limit: Max results

        Returns:
            List of images with relevance scores
        """
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()

        if document_name:
            cursor.execute("""
                SELECT id, image_path, document_name, page_number, figure_caption, context_text
                FROM images
                WHERE document_name = ? AND figure_caption IS NOT NULL AND figure_caption != ''
            """, (document_name,))
        else:
            cursor.execute("""
                SELECT id, image_path, document_name, page_number, figure_caption, context_text
                FROM images
                WHERE figure_caption IS NOT NULL AND figure_caption != ''
            """)

        rows = cursor.fetchall()
        conn.close()

        # Score each image by caption match
        scored_images = []
        query_lower = query.lower()
        query_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', query_lower))

        # Remove common stop words
        stop_words = {'what', 'is', 'the', 'a', 'an', 'how', 'does', 'show', 'me',
                      'explain', 'about', 'and', 'or', 'of', 'in', 'to', 'for', 'this', 'that'}
        query_words = query_words - stop_words

        for row in rows:
            img_id, img_path, doc_name, page_num, caption, context = row
            caption_lower = (caption or "").lower()

            # Calculate match score
            score = 0
            matched_words = []

            for word in query_words:
                if word in caption_lower:
                    score += len(word) * 2  # Double weight for caption matches
                    matched_words.append(word)

            if score > 0:
                scored_images.append({
                    "id": img_id,
                    "image_path": img_path,
                    "document_name": doc_name,
                    "page_number": page_num,
                    "figure_caption": caption,
                    "score": score,
                    "matched_words": matched_words,
                })

        # Sort by score descending
        scored_images.sort(key=lambda x: x["score"], reverse=True)

        return scored_images[:limit]

    def _search_by_context(self, query: str, document_name: str = None, limit: int = 5) -> List[Dict]:
        """
        Search images by matching query against layout-derived context.

        Args:
            query: Search query
            document_name: Optional document filter
            limit: Max results

        Returns:
            List of images with relevance scores
        """
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()

        if document_name:
            cursor.execute("""
                SELECT id, image_path, document_name, page_number, figure_caption, context_text
                FROM images
                WHERE document_name = ? AND context_text IS NOT NULL AND context_text != ''
            """, (document_name,))
        else:
            cursor.execute("""
                SELECT id, image_path, document_name, page_number, figure_caption, context_text
                FROM images
                WHERE context_text IS NOT NULL AND context_text != ''
            """)

        rows = cursor.fetchall()
        conn.close()

        # Score each image by context match
        scored_images = []
        query_lower = query.lower()
        query_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', query_lower))

        stop_words = {'what', 'is', 'the', 'a', 'an', 'how', 'does', 'show', 'me',
                      'explain', 'about', 'and', 'or', 'of', 'in', 'to', 'for', 'this', 'that'}
        query_words = query_words - stop_words

        for row in rows:
            img_id, img_path, doc_name, page_num, caption, context = row
            context_lower = (context or "").lower()

            score = 0
            for word in query_words:
                if word in context_lower:
                    score += len(word)

            if score > 0:
                scored_images.append({
                    "id": img_id,
                    "image_path": img_path,
                    "document_name": doc_name,
                    "page_number": page_num,
                    "figure_caption": caption,
                    "context_text": context,
                    "score": score,
                })

        scored_images.sort(key=lambda x: x["score"], reverse=True)

        return scored_images[:limit]

    def _extract_figure_references(self, text: str) -> List[str]:
        """
        Extract figure references from text (e.g., "Figure 3-1", "Fig. 5.2").

        Args:
            text: Text to search

        Returns:
            List of figure reference strings
        """
        patterns = [
            r'(?:Figure|Fig\.?)\s*(\d+[-.]?\d*)',
        ]

        references = []
        for pattern in patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                references.append(f"Figure {match}")

        return list(set(references))[:10]  # Dedupe and limit

    def _search_by_figure_reference(self, figure_ref: str, document_name: str = None) -> List[Dict]:
        """
        Search for images whose caption matches a figure reference.

        Args:
            figure_ref: Figure reference string (e.g., "Figure 3-1")
            document_name: Optional document filter

        Returns:
            List of matching images
        """
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()

        # Build search pattern
        fig_pattern = f"%{figure_ref}%"

        if document_name:
            cursor.execute("""
                SELECT id, image_path, document_name, page_number, figure_caption, context_text
                FROM images
                WHERE document_name = ?
                AND (figure_caption LIKE ? OR context_text LIKE ?)
                LIMIT 3
            """, (document_name, fig_pattern, fig_pattern))
        else:
            cursor.execute("""
                SELECT id, image_path, document_name, page_number, figure_caption, context_text
                FROM images
                WHERE figure_caption LIKE ? OR context_text LIKE ?
                LIMIT 3
            """, (fig_pattern, fig_pattern))

        rows = cursor.fetchall()
        conn.close()

        return [{
            "id": row[0],
            "image_path": row[1],
            "document_name": row[2],
            "page_number": row[3],
            "figure_caption": row[4],
            "context_text": row[5],
        } for row in rows]

    # =========================================================================
    # UTILITY METHODS
    # =========================================================================

    def get_images_for_pages(self, document_name: str, page_numbers: List[int], limit: int = 3) -> List[Dict]:
        """Get images from specific pages of a document (used as fallback only)"""
        if not page_numbers:
            return []

        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()

        placeholders = ",".join("?" * len(page_numbers))
        cursor.execute(f"""
            SELECT id, image_path, document_name, page_number, figure_caption
            FROM images
            WHERE document_name = ? AND page_number IN ({placeholders})
            ORDER BY page_number
            LIMIT ?
        """, [document_name] + page_numbers + [limit])

        rows = cursor.fetchall()
        conn.close()

        return [{
            "id": row[0],
            "image_path": row[1],
            "document_name": row[2],
            "page_number": row[3],
            "figure_caption": row[4],
        } for row in rows]

    def get_all_images(self) -> List[Dict]:
        """Get all stored images"""
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, image_path, document_name, page_number, figure_caption
            FROM images
        """)
        rows = cursor.fetchall()
        conn.close()

        return [{
            "id": row[0],
            "image_path": row[1],
            "document_name": row[2],
            "page_number": row[3],
            "figure_caption": row[4],
        } for row in rows]

    def get_images_without_caption(self) -> List[Dict]:
        """Get images that don't have figure captions"""
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, image_path, document_name, page_number
            FROM images
            WHERE figure_caption IS NULL OR figure_caption = ''
        """)
        rows = cursor.fetchall()
        conn.close()

        return [{
            "id": row[0],
            "image_path": row[1],
            "document_name": row[2],
            "page_number": row[3]
        } for row in rows]


# Global instance
image_extractor = None


def get_image_extractor() -> ImageExtractor:
    """Get or create image extractor singleton"""
    global image_extractor
    if image_extractor is None:
        image_extractor = ImageExtractor()
    return image_extractor
